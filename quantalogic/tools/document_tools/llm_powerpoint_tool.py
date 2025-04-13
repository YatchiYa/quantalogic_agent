"""Tool for generating well-structured PowerPoint presentations using LLM guidance."""

import json
import os
from pathlib import Path
from typing import Any, ClassVar, Dict, List, Optional, Tuple

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from quantalogic.tools.llm_tool import LLMTool
from quantalogic.tools.tool import Tool, ToolArgument


class LLMPowerPointTool(Tool):
    """Tool for generating professional PowerPoint presentations with LLM-guided structure."""

    name: str = "llm_powerpoint_tool"
    description: str = (
        "Creates well-structured PowerPoint presentations using LLM guidance. "
        "The LLM helps organize content into a professional presentation format "
        "with proper slide hierarchy, formatting, and visual balance. "
        "Saves output in /tmp directory."
    )
    need_validation: bool = False

    # Default style configuration
    DEFAULT_STYLES: ClassVar[Dict[str, Any]] = {
        "title_font": "Calibri",
        "body_font": "Calibri",
        "title_size": 44,
        "subtitle_size": 32,
        "heading_size": 28,
        "body_size": 18,
        "primary_color": [0, 112, 192],
        "secondary_color": [68, 114, 196],
        "text_color": [0, 0, 0],
        "background_color": [255, 255, 255],
    }

    # System prompt for the LLM to structure the presentation
    SYSTEM_PROMPT: ClassVar[str] = '''You are a professional presentation designer with expertise in creating clear, engaging, and well-structured PowerPoint presentations.

Your task is to organize the given content into a presentation outline with the following structure:

1. Title Slide
   - Main title
   - Optional subtitle
   - Any key introductory elements

2. Agenda/Overview
   - Clear outline of presentation sections
   - Key points to be covered

3. Content Slides (organized by sections)
   For each section:
   - Section title slide
   - Content slides with:
     * Clear headings
     * Bullet points (2-5 per slide)
     * Key messages
     * Supporting details
     * Visual suggestions (charts, diagrams, images)

4. Summary/Conclusion
   - Key takeaways
   - Call to action or next steps
   - Contact information if relevant

Guidelines:
- Keep slides focused and uncluttered
- Use consistent formatting
- Limit text per slide
- Suggest relevant visuals
- Maintain logical flow
- Use clear hierarchy

Output Format:
Return a JSON structure with this exact schema:
{
    "title": "Presentation Title",
    "subtitle": "Optional Subtitle",
    "slides": [
        {
            "type": "title|section|content|summary",
            "title": "Slide Title",
            "subtitle": "Optional Slide Subtitle",
            "content": [
                "Content point 1",
                "Content point 2"
            ],
            "notes": "Speaker notes and suggestions",
            "layout": "title|section|two_content|comparison|etc",
            "visual_suggestions": [
                "Suggestion for charts/images/diagrams"
            ]
        }
    ]
}'''

    arguments: List[ToolArgument] = [
        ToolArgument(
            name="content",
            arg_type="string",
            description="Content to be converted into a presentation",
            required=True,
            example='''# Project Proposal: AI Integration

## Overview
Our company needs to integrate AI capabilities to stay competitive.

## Current Challenges
- Manual data processing
- Slow customer response
- Limited analytics

## Proposed Solution
1. Implement ML models
2. Automate workflows
3. Add predictive analytics

## Benefits
- 50% faster processing
- 24/7 customer service
- Data-driven decisions

## Timeline & Budget
- 6 months implementation
- $500K investment
- ROI within 12 months''',
        ),
        ToolArgument(
            name="output_path",
            arg_type="string",
            description="Path for saving the PowerPoint file",
            required=True,
            example="/tmp/presentation.pptx",
        ),
        ToolArgument(
            name="model_name",
            arg_type="string",
            description="Name of the LLM model to use",
            required=False,
            default="gpt-4",
            example="gpt-4",
        ),
        ToolArgument(
            name="style_config",
            arg_type="string",
            description="JSON string with style settings",
            required=False,
            example='''{
    "title_font": "Calibri",
    "body_font": "Calibri",
    "title_size": 44,
    "subtitle_size": 32,
    "heading_size": 28,
    "body_size": 18,
    "primary_color": [0, 112, 192],
    "secondary_color": [68, 114, 196],
    "text_color": [0, 0, 0],
    "background_color": [255, 255, 255]
}''',
        ),
    ]

    def _normalize_path(self, path: str) -> Path:
        """Ensure output path is in /tmp directory.
        
        Args:
            path: Original path
            
        Returns:
            Normalized path in /tmp
        """
        if not path.startswith("/tmp/"):
            filename = os.path.basename(path)
            path = os.path.join("/tmp", filename)
        return Path(path).resolve()

    def _parse_style_config(self, style_config: Optional[str]) -> Dict[str, Any]:
        """Parse and validate style configuration.
        
        Args:
            style_config: JSON style configuration string
            
        Returns:
            Validated style configuration dict
        """
        try:
            if not style_config:
                return self.DEFAULT_STYLES.copy()
            
            config = json.loads(style_config)
            
            # Merge with defaults
            styles = self.DEFAULT_STYLES.copy()
            styles.update(config)
            
            return styles
            
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid style configuration JSON: {e}") from e
        except Exception as e:
            raise ValueError(f"Error parsing style configuration: {e}") from e

    def _get_layout(self, prs: Presentation, layout_type: str) -> int:
        """Get the appropriate slide layout index based on type.
        
        Args:
            prs: PowerPoint presentation
            layout_type: Desired layout type
            
        Returns:
            Layout index to use
        """
        # Map our layout types to PowerPoint's built-in layouts
        LAYOUT_MAPPING = {
            'title': 0,        # Title Slide
            'section': 1,      # Section Header
            'two_content': 3,  # Two Content
            'comparison': 4,   # Comparison
            'content': 5,      # Content with Caption
            'picture': 6,      # Picture with Caption
            'summary': 7,      # Blank Slide (we'll format it for summary)
        }
        return LAYOUT_MAPPING.get(layout_type, 5)  # Default to Content with Caption

    def _add_placeholder_content(self, placeholder, content: str, styles: Dict[str, Any]) -> None:
        """Add and format content in a placeholder shape.
        
        Args:
            placeholder: PowerPoint shape placeholder
            content: Text content to add
            styles: Style configuration
        """
        if not placeholder:
            return
            
        text_frame = placeholder.text_frame
        text_frame.clear()  # Clear existing text
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        
        run = p.add_run()
        run.text = content
        
        font = run.font
        font.name = styles['body_font']
        font.size = Pt(styles['body_size'])
        font.color.rgb = RGBColor(*styles['text_color'])

    def _create_bullet_points(self, shape, points: List[str], styles: Dict[str, Any], 
                            level: int = 0) -> None:
        """Create formatted bullet points in a shape.
        
        Args:
            shape: PowerPoint shape
            points: List of bullet points
            styles: Style configuration
            level: Indentation level
        """
        if not shape:
            return
            
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        for i, point in enumerate(points):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = point
            p.level = level
            p.alignment = PP_ALIGN.LEFT
            
            font = p.font
            font.name = styles['body_font']
            font.size = Pt(styles['body_size'])
            font.color.rgb = RGBColor(*styles['text_color'])

    def _create_slide(self, prs: Presentation, slide_data: Dict[str, Any], 
                     styles: Dict[str, Any]) -> None:
        """Create and populate a single slide.
        
        Args:
            prs: PowerPoint presentation
            slide_data: Slide content and layout data
            styles: Style configuration
        """
        # Get appropriate layout
        layout_idx = self._get_layout(prs, slide_data.get('layout', 'content'))
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Process title
        if slide.shapes.title:
            title_frame = slide.shapes.title.text_frame
            title_frame.clear()
            title_frame.word_wrap = True
            title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            p = title_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            
            run = p.add_run()
            run.text = slide_data['title']
            
            font = run.font
            font.name = styles['title_font']
            font.size = Pt(styles['title_size'] if slide_data['type'] == 'title' else styles['heading_size'])
            font.color.rgb = RGBColor(*styles['primary_color'])
            font.bold = True
        
        # Process subtitle if present
        if slide_data.get('subtitle'):
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # Subtitle
                    subtitle_frame = shape.text_frame
                    subtitle_frame.clear()
                    subtitle_frame.word_wrap = True
                    subtitle_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    
                    p = subtitle_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    
                    run = p.add_run()
                    run.text = slide_data['subtitle']
                    
                    font = run.font
                    font.name = styles['body_font']
                    font.size = Pt(styles['subtitle_size'])
                    font.color.rgb = RGBColor(*styles['secondary_color'])
                    break
        
        # Process content based on layout type
        if slide_data.get('content'):
            content_shapes = [shape for shape in slide.placeholders 
                            if shape.placeholder_format.type in [1, 3, 4]]  # Body text placeholders
            
            if content_shapes:
                if slide_data['layout'] == 'two_content' and len(content_shapes) >= 2:
                    # Split content between two columns
                    mid = len(slide_data['content']) // 2
                    self._create_bullet_points(content_shapes[0], slide_data['content'][:mid], styles)
                    self._create_bullet_points(content_shapes[1], slide_data['content'][mid:], styles)
                else:
                    # Use single content placeholder
                    self._create_bullet_points(content_shapes[0], slide_data['content'], styles)
        
        # Add notes if present
        if slide_data.get('notes'):
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = slide_data['notes']

    def execute(
        self,
        content: str,
        output_path: str,
        model_name: str = "gpt-4o-mini",
        style_config: Optional[str] = None,
    ) -> str:
        """Generate a PowerPoint presentation from the input content.

        Args:
            content: Input content to convert
            output_path: Where to save the PPTX file
            model_name: LLM model to use
            style_config: Optional style configuration

        Returns:
            Status message with output file information

        Raises:
            ValueError: If parameters are invalid
            Exception: For other conversion errors
        """
        try:
            # Initialize LLM tool
            llm = LLMTool(
                model_name=model_name,
                system_prompt=self.SYSTEM_PROMPT,
                name="presentation_designer"
            )
            
            # Get presentation structure from LLM
            llm_response = llm.execute(
                system_prompt=self.SYSTEM_PROMPT,
                prompt=content,
                temperature="0.7"
            )
            
            try:
                presentation_data = json.loads(llm_response)
            except json.JSONDecodeError as e:
                raise ValueError(f"Invalid LLM response format: {e}") from e
            
            # Parse style configuration
            styles = self._parse_style_config(style_config)
            
            # Create presentation
            prs = Presentation()
            
            # Set slide size to 16:9 (standard widescreen)
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
            # Create slides
            for slide_data in presentation_data['slides']:
                self._create_slide(prs, slide_data, styles)
            
            # Save presentation
            output_path = self._normalize_path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
            
            # Generate status message
            return (
                f"PowerPoint presentation created successfully:\n"
                f"- File: {output_path}\n"
                f"- Slides: {len(presentation_data['slides'])}\n"
                f"- Title: {presentation_data['title']}"
            )

        except Exception as e:
            logger.error(f"Error generating PowerPoint: {str(e)}")
            raise ValueError(f"Failed to generate PowerPoint: {str(e)}") from e


if __name__ == "__main__":
    # Example usage
    tool = LLMPowerPointTool()
    print(tool.to_markdown())
    
    # Test with sample content
    test_content = '''# Project Proposal: AI Integration

## Overview
Our company needs to integrate AI capabilities to stay competitive.

## Current Challenges
- Manual data processing
- Slow customer response
- Limited analytics

## Proposed Solution
1. Implement ML models
2. Automate workflows
3. Add predictive analytics

## Benefits
- 50% faster processing
- 24/7 customer service
- Data-driven decisions

## Timeline & Budget
- 6 months implementation
- $500K investment
- ROI within 12 months'''
    
    try:
        result = tool.execute(
            content=test_content,
            output_path="/tmp/test_presentation.pptx",
            model_name="gpt-4o-mini"
        )
        print("\nTest Result:")
        print(result)
    except Exception as e:
        print(f"Error: {e}")
